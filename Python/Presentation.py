from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Add cover slide
cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = cover_slide.shapes.title
subtitle = cover_slide.placeholders[1]
title.text = "Generators and Decorators in Python"
subtitle.text = "Author: Your Name\nDate: 2024-07-30"

# Add table of contents slide
toc_slide = prs.slides.add_slide(prs.slide_layouts[1])
toc_title = toc_slide.shapes.title
toc_content = toc_slide.placeholders[1]
toc_title.text = "Table of Contents"
toc_content.text = (
    "1. Introduction to Generators\n"
    "2. Generators Example\n"
    "3. Introduction to Decorators\n"
    "4. Decorators Example\n"
    "5. Conclusion"
)

# Add Introduction to Generators slide
gen_intro_slide = prs.slides.add_slide(prs.slide_layouts[1])
gen_intro_title = gen_intro_slide.shapes.title
gen_intro_content = gen_intro_slide.placeholders[1]
gen_intro_title.text = "Introduction to Generators"
gen_intro_content.text = (
    "Generators are a simple and powerful tool for creating iterators in Python. "
    "They are written like regular functions but use the 'yield' statement whenever they want to return data. "
    "Each time 'next()' is called on the generator, it resumes where it left off (it remembers all the data values and which statement was last executed). "
    "Key features of generators:\n"
    "- Memory Efficiency: Generators yield items one at a time instead of storing the entire sequence in memory.\n"
    "- Infinite Sequences: Generators can represent infinite sequences, such as streams of data from a sensor.\n"
    "- Pipeline Processing: Generators can be used to pipeline data processing tasks."
)

# Add Generators Example slide
gen_example_slide = prs.slides.add_slide(prs.slide_layouts[1])
gen_example_title = gen_example_slide.shapes.title
gen_example_content = gen_example_slide.placeholders[1]
gen_example_title.text = "Generators Example"
gen_example_content.text = (
    "Here is a simple example of a generator function in Python:\n\n"
    "def count_up_to(max):\n"
    "    count = 1\n"
    "    while count <= max:\n"
    "        yield count\n"
    "        count += 1\n\n"
    "# Using the generator\n"
    "counter = count_up_to(5)\n"
    "for number in counter:\n"
    "    print(number)\n\n"
    "This function generates numbers from 1 to a specified maximum value. "
    "Each time 'yield' is called, the current state of the function is saved, and the next number is produced when the function is resumed."
)

# Add Introduction to Decorators slide
dec_intro_slide = prs.slides.add_slide(prs.slide_layouts[1])
dec_intro_title = dec_intro_slide.shapes.title
dec_intro_content = dec_intro_slide.placeholders[1]
dec_intro_title.text = "Introduction to Decorators"
dec_intro_content.text = (
    "Decorators are a powerful and useful tool in Python that allow you to modify the behavior of a function or method. "
    "They are higher-order functions that take another function as an argument, extend its behavior, and return the modified function. "
    "Decorators are commonly used to add logging, validation, or other cross-cutting concerns to functions and methods.\n"
    "Key features of decorators:\n"
    "- Reusability: Decorators can be used to apply the same piece of logic to multiple functions, promoting code reusability.\n"
    "- Separation of Concerns: By using decorators, you can separate the core functionality of a function from its auxiliary functionalities."
)

# Add Decorators Example slide
dec_example_slide = prs.slides.add_slide(prs.slide_layouts[1])
dec_example_title = dec_example_slide.shapes.title
dec_example_content = dec_example_slide.placeholders[1]
dec_example_title.text = "Decorators Example"
dec_example_content.text = (
    "Here is a simple example of a decorator function in Python:\n\n"
    "def my_decorator(func):\n"
    "    def wrapper():\n"
    "        print('Something is happening before the function is called.')\n"
    "        func()\n"
    "        print('Something is happening after the function is called.')\n"
    "    return wrapper\n\n"
    "@my_decorator\n"
    "def say_hello():\n"
    "    print('Hello!')\n\n"
    "say_hello()\n\n"
    "This decorator adds some behavior before and after the 'say_hello' function is called."
)

# Add Conclusion slide
conclusion_slide = prs.slides.add_slide(prs.slide_layouts[1])
conclusion_title = conclusion_slide.shapes.title
conclusion_content = conclusion_slide.placeholders[1]
conclusion_title.text = "Conclusion"
conclusion_content.text = (
    "In this presentation, we covered two important features of Python: generators and decorators. "
    "Generators allow for efficient and lazy evaluation of sequences, making them memory efficient and suitable for representing infinite sequences. "
    "Decorators enable modification of function behavior in a reusable and maintainable way, promoting separation of concerns. "
    "Understanding and using these tools can greatly enhance your Python programming skills."
)
